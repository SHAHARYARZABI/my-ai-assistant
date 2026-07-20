from flask import Flask, request, jsonify, render_template
from openai import OpenAI
import os, uuid, json
from datetime import datetime, timezone
from dotenv import load_dotenv
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
import base64
from flask import send_from_directory
from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import re
import json as pyjson

load_dotenv()
print("MY KEY IS:", os.getenv("GROK_API_KEY"))

app = Flask(__name__)

# Grok uses OpenAI's SDK format, just pointed at xAI's servers
client = OpenAI(
    api_key=os.getenv("GROK_API_KEY"),
    base_url="https://api.groq.com/openai/v1"
)


CHATS_DIR = "chats"
os.makedirs(CHATS_DIR, exist_ok=True)
ASSISTANT_NAME = os.getenv("ASSISTANT_NAME", "My AI")

SYSTEM_PROMPT = f"You are {ASSISTANT_NAME}, a helpful, friendly personal AI assistant."

@app.route("/")
def home():
    return render_template("index.html", assistant_name=ASSISTANT_NAME)

@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.json
    chat_id = data.get("chat_id")
    user_message = data.get("message", "")
    attached_text = data.get("attached_text")
    image_data_url = data.get("image_data_url")

    if not chat_id:
        return jsonify({"error": "Missing chat_id"}), 400

    history = load_history(chat_id)

    if image_data_url:
        # Vision-style message: text + image together
        user_content = [
            {"type": "text", "text": user_message},
            {"type": "image_url", "image_url": {"url": image_data_url}}
        ]
        model_to_use = "llama-3.2-90b-vision-preview"  # confirm exact name in console.groq.com
    else:
        full_user_content = user_message
        if attached_text:
            full_user_content = f"{user_message}\n\n[Attached file content]:\n{attached_text[:6000]}"
        user_content = full_user_content
        model_to_use = "llama-3.3-70b-versatile"

    history["messages"].append({"role": "user", "content": user_content})
    api_messages = [{"role": "system", "content": SYSTEM_PROMPT}] + history["messages"]

    try:
        response = client.chat.completions.create(
            model=model_to_use,
            messages=api_messages,
            max_tokens=1200,
        )
        reply = response.choices[0].message.content
    except Exception as e:
        reply = f"Sorry, I hit an error reaching the model: {str(e)}"

    history["messages"].append({"role": "assistant", "content": reply})

    if history.get("title") in (None, "New chat") and user_message:
        history["title"] = (user_message[:40] + "…") if len(user_message) > 40 else user_message

    save_history(chat_id, history)
    return jsonify({"reply": reply, "title": history["title"]})

@app.route("/api/new_chat", methods=["POST"])
def new_chat():
    chat_id = str(uuid.uuid4())
    history = {"title": "New chat", "created": datetime.now(timezone.utc).isoformat(), "messages": []}
    save_history(chat_id, history)
    return jsonify({"chat_id": chat_id, "title": history["title"]})

@app.route("/api/list_chats")
def list_chats():
    chats = []
    for fname in os.listdir(CHATS_DIR):
        data = json.load(open(f"{CHATS_DIR}/{fname}"))
        chats.append({"chat_id": fname.replace(".json", ""), "title": data["title"], "created": data["created"]})
    chats.sort(key=lambda c: c["created"], reverse=True)
    return jsonify({"chats": chats})

@app.route("/api/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    filename = f.filename
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    text = ""
    is_image = False
    image_data_url = None

    try:
        if ext == "pdf":
            reader = PdfReader(f)
            pages_text = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(pages_text)

        elif ext == "docx":
            doc = Document(f)
            text = "\n".join(p.text for p in doc.paragraphs)

        elif ext == "pptx":
            prs = Presentation(f)
            slides_text = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                slide_text.append(run.text)
                slides_text.append(f"Slide {i}: " + " ".join(slide_text))
            text = "\n".join(slides_text)

        elif ext in ("txt", "csv", "md", "json", "py", "js", "html", "css"):
            raw = f.read()
            text = raw.decode("utf-8", errors="ignore")

        elif ext in ("jpg", "jpeg", "png", "webp", "gif"):
            is_image = True
            raw = f.read()
            b64 = base64.b64encode(raw).decode("utf-8")
            mime = f"image/{'jpeg' if ext == 'jpg' else ext}"
            image_data_url = f"data:{mime};base64,{b64}"
            text = "[Image attached — see image_data_url]"

        elif ext in ("mp4", "mov", "avi", "mkv", "mp3", "wav"):
            text = "[Video/audio file uploaded — this app cannot transcribe audio/video yet. Text-based questions about this file won't work until transcription is added.]"

        else:
            text = f"[Unsupported file type: .{ext}]"

    except Exception as e:
        text = f"[Could not read file: {str(e)}]"

    return jsonify({
        "filename": filename,
        "text": text[:8000],
        "is_image": is_image,
        "image_data_url": image_data_url
    })
    
GENERATED_DIR = "generated_files"
os.makedirs(GENERATED_DIR, exist_ok=True)


@app.route("/api/generate_file", methods=["POST"])
def generate_file():
    data = request.json
    topic = data.get("topic", "")
    file_type = data.get("file_type", "pptx")

    prompt = f"""Create content for a short presentation/document about: {topic}
Return ONLY valid JSON, no extra text, no markdown formatting, in exactly this format:
{{
  "title": "Main Title Here",
  "sections": [
    {{"heading": "Section heading", "bullets": ["point 1", "point 2", "point 3"]}}
  ]
}}
Include 5 to 7 sections."""

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2000,
        )
        raw = response.choices[0].message.content.strip()

        if raw.startswith("```"):
            raw = raw.strip("`")
            if raw.lower().startswith("json"):
                raw = raw[4:]

        match = re.search(r"\{.*\}", raw, re.DOTALL)
        content = pyjson.loads(match.group(0)) if match else pyjson.loads(raw)

    except Exception as e:
        return jsonify({"error": f"Could not generate content: {str(e)}"}), 500

    filename = f"{uuid.uuid4()}.{file_type}"
    filepath = os.path.join(GENERATED_DIR, filename)

    try:
        if file_type == "pptx":
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = content["title"]

            for sec in content["sections"]:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = sec["heading"]
                body = slide.placeholders[1].text_frame
                for i, bullet in enumerate(sec["bullets"]):
                    if i == 0:
                        body.text = bullet
                    else:
                        p = body.add_paragraph()
                        p.text = bullet
            prs.save(filepath)

        elif file_type == "docx":
            doc = Document()
            doc.add_heading(content["title"], level=0)
            for sec in content["sections"]:
                doc.add_heading(sec["heading"], level=1)
                for bullet in sec["bullets"]:
                    doc.add_paragraph(bullet, style="List Bullet")
            doc.save(filepath)

        elif file_type == "pdf":
            doc = SimpleDocTemplate(filepath, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [Paragraph(content["title"], styles["Title"]), Spacer(1, 20)]
            for sec in content["sections"]:
                story.append(Paragraph(sec["heading"], styles["Heading2"]))
                for bullet in sec["bullets"]:
                    story.append(Paragraph("• " + bullet, styles["Normal"]))
                story.append(Spacer(1, 12))
            doc.build(story)

        else:
            return jsonify({"error": "Unsupported file type"}), 400

    except Exception as e:
        return jsonify({"error": f"Could not build file: {str(e)}"}), 500

    return jsonify({"download_url": f"/files/{filename}", "filename": filename})


@app.route("/files/<filename>")
def download_generated_file(filename):
    return send_from_directory(GENERATED_DIR, filename, as_attachment=True)


def load_history(chat_id):
    path = f"{CHATS_DIR}/{chat_id}.json"
    return json.load(open(path)) if os.path.exists(path) else {"title": "New chat", "messages": []}

def save_history(chat_id, history):
    json.dump(history, open(f"{CHATS_DIR}/{chat_id}.json", "w"))

if __name__ == "__main__":
    app.run(debug=True, port=5000)